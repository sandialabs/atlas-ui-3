"""
PowerPoint Generator MCP Server using FastMCP.

Converts markdown content into a professional PowerPoint presentation with 16:9 aspect ratio.
Markdown headers (# or ##) become slide titles and content below each header becomes slide content.
Supports bullet point lists and optional image integration.

Tools:
 - markdown_to_pptx: Converts markdown content to PowerPoint presentation

Demonstrates: Markdown parsing, file output with base64 encoding, and professional templating.
"""

from __future__ import annotations

import base64
import html
import io
import logging
import os
import re
import tempfile
from pathlib import Path
from typing import Annotated, Any, Dict, List, Optional

import requests
from fastmcp import FastMCP
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Configuration
VERBOSE = True

# Configure logging first
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - PPTX_GENERATOR - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Configure paths and add file handler if available
current_dir = Path(__file__).parent
logger.info(f"Current dir: {current_dir.absolute()}")
backend_dir = current_dir.parent.parent
logger.info(f"Backend dir: {backend_dir.absolute()}")
project_root = backend_dir.parent
logger.info(f"Project root: {project_root.absolute()}")
logs_dir = project_root / 'logs'
logger.info(f"Logs dir: {logs_dir.absolute()}")
main_log_path = logs_dir / 'app.jsonl'
logger.info(f"Log path: {main_log_path.absolute()}")
logger.info(f"Log path exists: {main_log_path.exists()}")
logger.info(f"Logs dir exists: {logs_dir.exists()}")

# Add file handler if log path exists
if main_log_path.exists():
    file_handler = logging.FileHandler(main_log_path)
    file_handler.setFormatter(logging.Formatter('%(asctime)s - PPTX_GENERATOR - %(name)s - %(levelname)s - %(message)s'))
    logger.addHandler(file_handler)

mcp = FastMCP("pptx_generator")

# Sandia National Laboratories color scheme
SANDIA_BLUE = RGBColor(0, 51, 102)  # Dark blue - primary brand color
SANDIA_LIGHT_BLUE = RGBColor(0, 102, 153)  # Lighter blue for accents
SANDIA_RED = RGBColor(153, 0, 0)  # Red accent
SANDIA_GRAY = RGBColor(102, 102, 102)  # Gray for secondary text
SANDIA_WHITE = RGBColor(255, 255, 255)

# 16:9 slide dimensions (standard widescreen)
# Height is 7.5 inches, width is calculated for exact 16:9 ratio
SLIDE_HEIGHT = Inches(7.5)
SLIDE_WIDTH = Inches(7.5 * 16 / 9)  # 16:9 aspect ratio = 13.333... inches

# Allowed base paths for local file access (security constraint)
# Only allow files relative to the current working directory
ALLOWED_BASE_PATH = Path(".").resolve()


def _escape_html(text: str) -> str:
    """Escape HTML special characters to prevent XSS attacks."""
    return html.escape(text, quote=True)


def _is_safe_local_path(filepath: str) -> bool:
    """Check if a local file path is safe (within allowed base directory).

    Prevents path traversal attacks by ensuring the resolved path
    is within the allowed base directory.

    Args:
        filepath: The file path to validate

    Returns:
        True if the path is safe to access, False otherwise
    """
    if not filepath:
        return False

    try:
        requested_path = Path(filepath)
        if requested_path.is_absolute():
            resolved_path = requested_path.resolve()
        else:
            resolved_path = (ALLOWED_BASE_PATH / requested_path).resolve()

        # Ensure the path is within the allowed base directory
        resolved_path.relative_to(ALLOWED_BASE_PATH)
        return True
    except (ValueError, OSError):
        return False


def _calculate_indent_level(leading_spaces: int) -> int:
    """Calculate bullet indent level from leading whitespace."""
    return leading_spaces // 2 if leading_spaces >= 2 else 0


def _clean_markdown_text(text: str) -> str:
    """Clean markdown formatting from text while preserving content."""
    if not text:
        return ""

    # Remove bold markers (**text** or __text__)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'__(.+?)__', r'\1', text)

    # Remove italic markers (*text* or _text_) - be careful not to match bullet points
    text = re.sub(r'(?<!\*)\*([^*\n]+?)\*(?!\*)', r'\1', text)
    text = re.sub(r'(?<!_)_([^_\n]+?)_(?!_)', r'\1', text)

    # Remove inline code markers (`text`)
    text = re.sub(r'`([^`]+?)`', r'\1', text)

    # Remove image syntax ![alt](url) - must come before link syntax
    text = re.sub(r'!\[([^\]]*?)\]\([^)]+?\)', r'\1', text)

    # Remove link syntax [text](url) - keep the text
    text = re.sub(r'\[([^\]]+?)\]\([^)]+?\)', r'\1', text)

    # Clean up any remaining markdown artifacts
    text = re.sub(r'^\s*#{1,6}\s*', '', text)  # Remove header markers at start of line

    return text.strip()


def _apply_sandia_template(prs: Presentation) -> None:
    """Apply Sandia-style template settings to the presentation."""
    # Set 16:9 aspect ratio
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT


def _add_footer_bar(slide_obj, slide_num: int, total_slides: int) -> None:
    """Add a professional footer bar to the slide."""
    # Add footer bar at bottom
    footer_height = Inches(0.4)
    footer_top = SLIDE_HEIGHT - footer_height

    # Add footer rectangle
    footer_shape = slide_obj.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), footer_top,
        SLIDE_WIDTH, footer_height
    )
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = SANDIA_BLUE
    footer_shape.line.fill.background()

    # Add slide number text
    slide_num_box = slide_obj.shapes.add_textbox(
        SLIDE_WIDTH - Inches(1), footer_top + Inches(0.05),
        Inches(0.9), Inches(0.3)
    )
    tf = slide_num_box.text_frame
    tf.paragraphs[0].text = f"{slide_num}/{total_slides}"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = SANDIA_WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _add_header_bar(slide_obj) -> None:
    """Add a professional header accent bar to the slide."""
    # Add thin accent bar at top
    header_height = Inches(0.1)
    header_shape = slide_obj.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, header_height
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = SANDIA_LIGHT_BLUE
    header_shape.line.fill.background()


def _style_title(title_shape) -> None:
    """Apply Sandia styling to title text."""
    if title_shape and title_shape.has_text_frame:
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = SANDIA_BLUE
            paragraph.font.bold = True
            paragraph.font.size = Pt(36)

def _sanitize_filename(filename: str, max_length: int = 50) -> str:
    """Sanitize filename by removing bad characters and truncating."""
    # Remove bad characters (anything not alphanumeric, underscore, or dash)
    cleaned_filename = re.sub(r'[^\w\-]', '', filename)
    # Remove newlines and extra spaces
    cleaned_filename = re.sub(r'\s+', '', cleaned_filename)
    # Truncate to max length
    return cleaned_filename[:max_length] if cleaned_filename else "presentation"

def _is_backend_download_path(s: str) -> bool:
    """Detect backend-relative download paths like /api/files/download/...."""
    return isinstance(s, str) and s.startswith("/api/files/download/")


def _backend_base_url() -> str:
    """Resolve backend base URL from environment variable."""
    return os.environ.get("CHATUI_BACKEND_BASE_URL", "http://127.0.0.1:8000")


def _load_image_bytes(filename: str, file_data_base64: str = "") -> Optional[bytes]:
    """Load image data from filename or base64 data."""
    if file_data_base64:
        try:
            return base64.b64decode(file_data_base64)
        except Exception as e:
            if VERBOSE:
                logger.info(f"Error decoding base64 image data: {e}")
            return None

    if _is_backend_download_path(filename):
        # Backend provided a download path
        full_url = _backend_base_url() + filename
        try:
            if VERBOSE:
                logger.info(f"Fetching image from {full_url}")
            response = requests.get(full_url, timeout=30)
            response.raise_for_status()
            return response.content
        except Exception as e:
            if VERBOSE:
                logger.info(f"Error fetching image from {full_url}: {e}")
            return None

    # Try as local file path - with path traversal protection
    if _is_safe_local_path(filename) and os.path.isfile(filename):
        try:
            with open(filename, "rb") as f:
                return f.read()
        except Exception as e:
            if VERBOSE:
                logger.info(f"Error reading local image file {filename}: {e}")
            return None
    elif not _is_safe_local_path(filename):
        if VERBOSE:
            logger.warning(f"Blocked access to unsafe path: {filename}")
        return None

    if VERBOSE:
        logger.info(f"Image file not found: {filename}")
    return None


def _parse_markdown_slides(markdown_content: str) -> List[Dict[str, str]]:
    """Parse markdown content into slides with improved bullet point handling."""
    slides = []

    # Split by headers (# or ##)
    sections = re.split(r'^#{1,2}\s+(.+)$', markdown_content, flags=re.MULTILINE)

    # Remove empty first element if exists
    if sections and not sections[0].strip():
        sections = sections[1:]

    # Group into title/content pairs
    for i in range(0, len(sections), 2):
        if i + 1 < len(sections):
            title = _clean_markdown_text(sections[i].strip())
            content = sections[i + 1].strip() if i + 1 < len(sections) else ""
            slides.append({"title": title, "content": content})
        elif sections[i].strip():
            # Handle case where there's a title but no content
            slides.append({"title": _clean_markdown_text(sections[i].strip()), "content": ""})

    # If no headers found, treat entire content as one slide
    if not slides and markdown_content.strip():
        slides.append({"title": "Slide 1", "content": markdown_content.strip()})

    return slides


def _add_image_to_slide(slide_obj, image_bytes: bytes,
                       left: Inches = Inches(1), top: Inches = Inches(2),
                       width: Inches = Inches(10), height: Inches = Inches(5)):
    """Add image to a slide with 16:9 optimized positioning."""
    try:
        # Create a temporary file for the image
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            tmp_file.write(image_bytes)
            tmp_file.flush()

            # Add image to slide
            pic = slide_obj.shapes.add_picture(tmp_file.name, left, top, width, height)

            # Clean up
            os.unlink(tmp_file.name)

            return pic
    except Exception as e:
        logger.error(f"Error adding image to slide: {e}")
        return None



@mcp.tool
def markdown_to_pptx(
    markdown_content: Annotated[str, "Markdown content with headers (# or ##) as slide titles and content below each header"],
    file_name: Annotated[Optional[str], "Output file name (base name for generated files without extension)"] = None,
    image_filename: Annotated[Optional[str], "Optional image filename to integrate into the presentation"] = "",
    image_data_base64: Annotated[Optional[str], "Framework may supply Base64 image content as fallback"] = ""
) -> Dict[str, Any]:
    """
    Converts markdown content to a professional PowerPoint presentation with 16:9 aspect ratio.

    Creates polished presentations with Sandia-style professional templating, proper bullet
    point formatting, and clean markdown text handling.

    Args:
        markdown_content: Markdown content where headers (# or ##) become slide titles and content below becomes slide content
        file_name: Output file name (base name for generated files without extension)
        image_filename: Optional image filename to integrate into the presentation
        image_data_base64: Framework may supply Base64 image content as fallback

    Returns:
        Dictionary with 'results' and 'artifacts' keys:
        - 'results': Success message or error message
        - 'artifacts': List of artifact dictionaries with 'name', 'b64', and 'mime' keys
    """
    if VERBOSE:
        logger.info("Starting markdown_to_pptx execution...")
    try:
        # Handle None values and sanitize the output filename
        image_filename = image_filename or ""
        image_data_base64 = image_data_base64 or ""
        # Use file_name if provided, otherwise use default "presentation"
        output_filename = _sanitize_filename(file_name or "presentation")

        # Parse markdown into slides
        slides = _parse_markdown_slides(markdown_content)
        if VERBOSE:
            logger.info(f"Parsed {len(slides)} slides from markdown")

        if not slides:
            return {"results": {"error": "No slides could be parsed from markdown content"}}

        total_slides = len(slides)

        # Load image if provided
        image_bytes = None
        if image_filename:
            image_bytes = _load_image_bytes(image_filename, image_data_base64)
            if image_bytes:
                if VERBOSE:
                    logger.info(f"Loaded image: {image_filename}")
            else:
                if VERBOSE:
                    logger.info(f"Failed to load image: {image_filename}")

        # Create presentation with 16:9 aspect ratio
        prs = Presentation()
        _apply_sandia_template(prs)
        if VERBOSE:
            logger.info("Created PowerPoint presentation with 16:9 aspect ratio")

        for i, slide_data in enumerate(slides):
            title = slide_data.get('title', 'Untitled Slide')
            content = slide_data.get('content', '')

            # Add slide using blank layout for full control
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide_obj = prs.slides.add_slide(slide_layout)

            # Add header accent bar
            _add_header_bar(slide_obj)

            # Add title text box
            title_box = slide_obj.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                SLIDE_WIDTH - Inches(1), Inches(0.8)
            )
            title_tf = title_box.text_frame
            title_tf.word_wrap = True
            title_p = title_tf.paragraphs[0]
            title_p.text = title
            title_p.font.size = Pt(36)
            title_p.font.bold = True
            title_p.font.color.rgb = SANDIA_BLUE
            title_p.alignment = PP_ALIGN.LEFT

            if VERBOSE:
                logger.info(f"Added slide {i+1}: {title}")

            # Add content text box
            content_box = slide_obj.shapes.add_textbox(
                Inches(0.5), Inches(1.3),
                SLIDE_WIDTH - Inches(1), SLIDE_HEIGHT - Inches(2.0)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            # Process content - handle bullet points and regular text with improved cleanup
            if content.strip():
                lines = content.split('\n')
                first_paragraph = True

                for line in lines:
                    line = line.rstrip()

                    if not line.strip():
                        continue

                    # Calculate indentation level
                    indent_level = 0
                    stripped_line = line.lstrip()
                    leading_spaces = len(line) - len(stripped_line)

                    # Check for various bullet point formats
                    is_bullet = False
                    bullet_text = stripped_line

                    # Handle numbered lists (1. 2. etc.)
                    numbered_match = re.match(r'^(\d+)\.\s+(.+)$', stripped_line)
                    if numbered_match:
                        is_bullet = True
                        bullet_text = numbered_match.group(2)
                        indent_level = _calculate_indent_level(leading_spaces)
                    # Handle bullet points (-, *, +) with regex for proper text extraction
                    else:
                        bullet_match = re.match(r'^[-*+]\s+(.+)$', stripped_line)
                        if bullet_match:
                            is_bullet = True
                            bullet_text = bullet_match.group(1)
                            indent_level = _calculate_indent_level(leading_spaces)

                    # Clean the bullet text from markdown formatting
                    bullet_text = _clean_markdown_text(bullet_text.strip())

                    if not bullet_text:
                        continue

                    if first_paragraph:
                        p = tf.paragraphs[0]
                        first_paragraph = False
                    else:
                        p = tf.add_paragraph()

                    p.text = bullet_text
                    p.level = min(indent_level, 4)  # Cap at level 4
                    p.font.size = Pt(20)
                    p.font.color.rgb = SANDIA_GRAY if indent_level > 0 else RGBColor(51, 51, 51)
                    p.space_after = Pt(8)
                    p.alignment = PP_ALIGN.LEFT

            # Add footer bar with slide number
            _add_footer_bar(slide_obj, i + 1, total_slides)

            # Add image to first slide if provided
            if i == 0 and image_bytes:
                _add_image_to_slide(slide_obj, image_bytes,
                                   left=Inches(8), top=Inches(2),
                                   width=Inches(4.5), height=Inches(4))

        # Write outputs to a temporary directory and clean up after encoding
        with tempfile.TemporaryDirectory() as tmpdir:
            # Save presentation
            pptx_output_path = os.path.join(tmpdir, f"output_{output_filename}.pptx")
            prs.save(pptx_output_path)
            if VERBOSE:
                logger.info(f"Saved PowerPoint presentation to {pptx_output_path}")

            # Create HTML file instead of PDF
            html_output_path = os.path.join(tmpdir, f"output_{output_filename}.html")
            if VERBOSE:
                logger.info(f"Starting HTML creation to {html_output_path}")

            # Create HTML representation of the presentation with Sandia styling
            html_content = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PowerPoint Presentation</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; background-color: #f0f0f0; }
        .slide {
            background: white;
            margin: 20px auto;
            padding: 0;
            max-width: 960px;
            aspect-ratio: 16/9;
            box-shadow: 0 4px 12px rgba(0,0,0,0.15);
            border-radius: 4px;
            page-break-after: always;
            position: relative;
            overflow: hidden;
        }
        .header-bar {
            background: #006699;
            height: 8px;
            width: 100%;
        }
        .footer-bar {
            background: #003366;
            height: 32px;
            width: 100%;
            position: absolute;
            bottom: 0;
            display: flex;
            align-items: center;
            justify-content: flex-end;
            padding-right: 20px;
            box-sizing: border-box;
        }
        .slide-number {
            color: white;
            font-size: 12px;
        }
        .slide-content-wrapper {
            padding: 20px 40px;
        }
        .slide-title {
            color: #003366;
            font-size: 28px;
            font-weight: bold;
            margin-bottom: 20px;
        }
        .slide-content {
            font-size: 16px;
            line-height: 1.6;
            color: #333;
        }
        .slide-content ul {
            margin: 0;
            padding-left: 30px;
            list-style-type: disc;
        }
        .slide-content ul ul {
            list-style-type: circle;
            color: #666;
        }
        .slide-content li {
            margin-bottom: 8px;
        }
        .slide-image {
            max-width: 350px;
            max-height: 250px;
            display: block;
            margin: 20px auto;
            border-radius: 4px;
        }
    </style>
</head>
<body>"""

            for i, slide_data in enumerate(slides):
                title = slide_data.get('title', 'Untitled Slide')
                content = slide_data.get('content', '')

                # Escape HTML to prevent XSS attacks
                safe_title = _escape_html(_clean_markdown_text(title))

                html_content += f"""
    <div class="slide">
        <div class="header-bar"></div>
        <div class="slide-content-wrapper">
            <div class="slide-title">{safe_title}</div>
            <div class="slide-content">"""

                # Add image to first slide if provided
                if i == 0 and image_bytes:
                    try:
                        img = Image.open(io.BytesIO(image_bytes))
                        mime_type = Image.MIME.get(img.format)
                        if mime_type:
                            img_b64 = base64.b64encode(image_bytes).decode('utf-8')
                            html_content += f'<img src="data:{mime_type};base64,{img_b64}" class="slide-image" />'
                    except Exception as e:
                        if VERBOSE:
                            logger.warning(f"Could not process image for HTML conversion: {e}")

                if content.strip():
                    lines = content.split('\n')
                    in_list = False
                    current_indent = 0

                    for line in lines:
                        stripped = line.strip()
                        if not stripped:
                            continue

                        # Detect bullet points with various formats
                        is_bullet = False
                        bullet_text = stripped
                        indent = _calculate_indent_level(len(line) - len(line.lstrip()))

                        # Numbered list
                        numbered_match = re.match(r'^(\d+)\.\s+(.+)$', stripped)
                        if numbered_match:
                            is_bullet = True
                            bullet_text = numbered_match.group(2)
                        else:
                            # Handle bullet points (-, *, +) with regex
                            bullet_match = re.match(r'^[-*+]\s+(.+)$', stripped)
                            if bullet_match:
                                is_bullet = True
                                bullet_text = bullet_match.group(1)

                        # Clean markdown from text and escape HTML
                        bullet_text = _escape_html(_clean_markdown_text(bullet_text))

                        if is_bullet:
                            if not in_list:
                                html_content += "<ul>"
                                in_list = True
                                current_indent = indent

                            # Handle indent changes
                            while current_indent < indent:
                                html_content += "<ul>"
                                current_indent += 1
                            while current_indent > indent:
                                html_content += "</ul>"
                                current_indent -= 1

                            html_content += f"<li>{bullet_text}</li>"
                        else:
                            # Close all open lists
                            while current_indent > 0:
                                html_content += "</ul>"
                                current_indent -= 1
                            if in_list:
                                html_content += "</ul>"
                                in_list = False
                            # Escape HTML in paragraph text to prevent XSS
                            safe_paragraph = _escape_html(_clean_markdown_text(stripped))
                            html_content += f"<p>{safe_paragraph}</p>"

                    # Close any remaining open lists
                    while current_indent > 0:
                        html_content += "</ul>"
                        current_indent -= 1
                    if in_list:
                        html_content += "</ul>"

                html_content += f"""
            </div>
        </div>
        <div class="footer-bar">
            <span class="slide-number">{i+1}/{total_slides}</span>
        </div>
    </div>"""

            html_content += """
</body>
</html>"""

            # Save HTML file
            try:
                if VERBOSE:
                    logger.info("Saving HTML file...")
                with open(html_output_path, "w", encoding="utf-8") as f:
                    f.write(html_content)
                if VERBOSE:
                    logger.info(f"HTML successfully created at {html_output_path}")
            except Exception as e:
                # If HTML save fails, continue with just PPTX
                if VERBOSE:
                    logger.warning(f"HTML creation failed: {str(e)}")
                # Remove the HTML file if it exists
                if os.path.exists(html_output_path):
                    os.remove(html_output_path)
                if VERBOSE:
                    logger.info("HTML file removed due to creation error")

            # Read PPTX file as bytes
            with open(pptx_output_path, "rb") as f:
                pptx_bytes = f.read()

            # Encode PPTX as base64
            pptx_b64 = base64.b64encode(pptx_bytes).decode('utf-8')
            if VERBOSE:
                logger.info("PPTX file successfully encoded to base64")

            # Read HTML file as bytes if it exists
            html_b64 = None
            if os.path.exists(html_output_path):
                with open(html_output_path, "r", encoding="utf-8") as f:
                    html_content_file = f.read()
                html_b64 = base64.b64encode(html_content_file.encode('utf-8')).decode('utf-8')
                if VERBOSE:
                    logger.info("HTML file successfully encoded to base64")

            # Prepare artifacts
            artifacts = [
                {
                    "name": f"{output_filename}.pptx",
                    "b64": pptx_b64,
                    "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                }
            ]

            # Add HTML if creation was successful
            if html_b64:
                artifacts.append({
                    "name": f"{output_filename}.html",
                    "b64": html_b64,
                    "mime": "text/html",
                })
                if VERBOSE:
                    logger.info(f"Added {len(artifacts)} artifacts to response")
            else:
                if VERBOSE:
                    logger.info("No HTML artifact added due to creation failure")

            return {
                "results": {
                    "operation": "markdown_to_pptx",
                    "message": "PowerPoint presentation and HTML file generated successfully from markdown.",
                    "html_generated": html_b64 is not None,
                    "image_included": image_bytes is not None,
                },
                "artifacts": artifacts,
                "display": {
                    "open_canvas": True,
                    "primary_file": f"{output_filename}.pptx",
                    "mode": "replace",
                    "viewer_hint": "powerpoint",
                },
                "meta_data": {
                    "generated_slides": len(slides),
                    "output_files": [f"{output_filename}.pptx", f"{output_filename}.html"] if html_b64 else [f"{output_filename}.pptx"],
                    "output_file_paths": [f"temp:output_{output_filename}.pptx", f"temp:output_{output_filename}.html"] if html_b64 else [f"temp:output_{output_filename}.pptx"],
                },
            }
    except Exception as e:
        logger.error(f"Error in markdown_to_pptx: {str(e)}")
        return {"results": {"error": f"Error creating PowerPoint from markdown: {str(e)}"}}


if __name__ == "__main__":
    mcp.run()
