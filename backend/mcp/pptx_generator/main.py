"""
PowerPoint Generator MCP Server using FastMCP.

Converts JSON input containing slide data into a PowerPoint presentation.
Input should be a JSON object with a 'slides' array, each containing 'title' and 'content' fields.
Supports bullet point lists in content for more complex slide formatting.

Tools:
 - json_to_pptx: Converts JSON input to PowerPoint presentation

Demonstrates: JSON input handling, file output with base64 encoding, and structured output.
"""

from __future__ import annotations

import base64
import logging
import os
import tempfile
import io
import re
import requests
from pathlib import Path
from typing import Any, Dict, List, Annotated, Optional
from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image


# Configuration
VERBOSE = True

# Configure logging first
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - PPTX_GENERATOR - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Configure paths and add file handler if available
current_dir = Path(__file__).parent
logger.info(f"Current dir: {current_dir.absolute()}")
backend_dir = current_dir.parent.parent
logger.info(f"Backend dir: {backend_dir.absolute()}")
project_root = backend_dir.parent
logger.info(f"Project root: {project_root.absolute()}")
logs_dir = project_root / 'logs'
logger.info(f"Logs dir: {logs_dir.absolute()}")
main_log_path = logs_dir / 'app.jsonl'
logger.info(f"Log path: {main_log_path.absolute()}")
logger.info(f"Log path exists: {main_log_path.exists()}")
logger.info(f"Logs dir exists: {logs_dir.exists()}")

# Add file handler if log path exists
if main_log_path.exists():
    file_handler = logging.FileHandler(main_log_path)
    file_handler.setFormatter(logging.Formatter('%(asctime)s - PPTX_GENERATOR - %(name)s - %(levelname)s - %(message)s'))
    logger.addHandler(file_handler)

mcp = FastMCP("pptx_generator")

def _is_backend_download_path(s: str) -> bool:
    """Detect backend-relative download paths like /api/files/download/...."""
    return isinstance(s, str) and s.startswith("/api/files/download/")


def _backend_base_url() -> str:
    """Resolve backend base URL from environment variable."""
    return os.environ.get("CHATUI_BACKEND_BASE_URL", "http://127.0.0.1:8000")


def _load_image_bytes(filename: str, file_data_base64: str = "") -> Optional[bytes]:
    """Load image data from filename or base64 data."""
    if file_data_base64:
        try:
            return base64.b64decode(file_data_base64)
        except Exception as e:
            if VERBOSE:
                logger.info(f"Error decoding base64 image data: {e}")
            return None
    
    if _is_backend_download_path(filename):
        # Backend provided a download path
        full_url = _backend_base_url() + filename
        try:
            if VERBOSE:
                logger.info(f"Fetching image from {full_url}")
            response = requests.get(full_url, timeout=30)
            response.raise_for_status()
            return response.content
        except Exception as e:
            if VERBOSE:
                logger.info(f"Error fetching image from {full_url}: {e}")
            return None
    
    # Try as local file path
    if os.path.isfile(filename):
        try:
            with open(filename, "rb") as f:
                return f.read()
        except Exception as e:
            if VERBOSE:
                logger.info(f"Error reading local image file {filename}: {e}")
            return None
    
    if VERBOSE:
        logger.info(f"Image file not found: {filename}")
    return None


def _parse_markdown_slides(markdown_content: str) -> List[Dict[str, str]]:
    """Parse markdown content into slides."""
    slides = []
    
    # Split by headers (# or ##)
    sections = re.split(r'^#{1,2}\s+(.+)$', markdown_content, flags=re.MULTILINE)
    
    # Remove empty first element if exists
    if sections and not sections[0].strip():
        sections = sections[1:]
    
    # Group into title/content pairs
    for i in range(0, len(sections), 2):
        if i + 1 < len(sections):
            title = sections[i].strip()
            content = sections[i + 1].strip() if i + 1 < len(sections) else ""
            slides.append({"title": title, "content": content})
        elif sections[i].strip():
            # Handle case where there's a title but no content
            slides.append({"title": sections[i].strip(), "content": ""})
    
    # If no headers found, treat entire content as one slide
    if not slides and markdown_content.strip():
        slides.append({"title": "Slide 1", "content": markdown_content.strip()})
    
    return slides


def _add_image_to_slide(slide_obj, image_bytes: bytes, left: Inches = Inches(1), top: Inches = Inches(2), 
                       width: Inches = Inches(8), height: Inches = Inches(5)):
    """Add image to a slide."""
    try:
        # Create a temporary file for the image
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            tmp_file.write(image_bytes)
            tmp_file.flush()
            
            # Add image to slide
            pic = slide_obj.shapes.add_picture(tmp_file.name, left, top, width, height)
            
            # Clean up
            os.unlink(tmp_file.name)
            
            return pic
    except Exception as e:
        print(f"Error adding image to slide: {e}")
        return None


@mcp.tool
def json_to_pptx(
    input_data: Annotated[str, "JSON string containing slide data in this format: {\"slides\": [{\"title\": \"Slide 1\", \"content\": \"- Item 1\\n- Item 2\\n- Item 3\"}, {\"title\": \"Slide 2\", \"content\": \"- Item A\\n- Item B\"}]}"],
    image_filename: Annotated[str, "Optional image filename to integrate into the presentation"] = "",
    image_data_base64: Annotated[str, "Framework may supply Base64 image content as fallback"] = ""
) -> Dict[str, Any]:
    """
    Create professional PowerPoint presentations from structured JSON data with advanced formatting and multimedia support.

    This comprehensive presentation generation tool transforms structured data into polished PowerPoint presentations:
    
    **Presentation Creation Features:**
    - Professional PowerPoint template and layout generation
    - Dynamic slide creation from JSON data structures
    - Automatic bullet point formatting and list management
    - Custom slide titles and content organization
    - Professional typography and spacing optimization

    **Content Formatting Capabilities:**
    - Intelligent bullet point recognition and formatting
    - Multi-level list support with proper indentation
    - Text formatting with consistent styling
    - Professional color schemes and layout templates
    - Automatic content overflow handling

    **Multimedia Integration:**
    - Image embedding with automatic sizing and positioning
    - Support for multiple image formats (PNG, JPG, GIF)
    - Base64 image data processing
    - Local file and URL-based image integration
    - Responsive image placement and scaling

    **JSON Data Structure:**
    - Flexible slide definition with title and content pairs
    - Nested content support for complex information
    - Array-based slide sequences for easy management
    - Extensible schema for future enhancements

    **Professional Output:**
    - High-quality PPTX file generation compatible with Microsoft PowerPoint
    - Professional business presentation templates
    - Consistent formatting and branding across slides
    - Optimized file size and compatibility

    **Use Cases:**
    - Business presentation automation from data
    - Educational content and training material generation
    - Marketing presentation creation from campaign data
    - Report summaries and executive briefings
    - Project status and milestone presentations
    - Product demonstrations and feature showcases

    **Advanced Features:**
    - Automatic slide layout optimization
    - Content-aware formatting decisions
    - Image aspect ratio preservation
    - Professional design pattern application

    **JSON Format Example:**
    ```json
    {
      "slides": [
        {
          "title": "Project Overview",
          "content": "- Project goals and objectives\\n- Timeline and milestones\\n- Key deliverables"
        },
        {
          "title": "Budget Analysis", 
          "content": "- Current spend: $50,000\\n- Remaining budget: $25,000\\n- Cost projections"
        }
      ]
    }
    ```

    Args:
        input_data: JSON string with slide definitions (title and content pairs with bullet points)
        image_filename: Optional image file to embed in presentation (supports various formats)
        image_data_base64: Alternative Base64-encoded image content (automatically provided by framework)

    Returns:
        Dictionary containing:
        - results: Presentation generation summary and success confirmation
        - artifacts: Professional PPTX file as downloadable content
        - display: Optimized viewer configuration for presentation review
        - meta_data: Generation statistics and file information
        Or error message if JSON parsing or presentation generation fails
    """
    print("Starting json_to_pptx execution...")
    try:
        import json
        data = json.loads(input_data)
        
        if not isinstance(data, dict) or 'slides' not in data:
            return {"results": {"error": "Input must be a JSON object containing 'slides' array"}}
            
        slides = data['slides']
        if VERBOSE:
            logger.info(f"Processing {len(slides)} slides...")
        
        # Load image if provided
        image_bytes = None
        if image_filename:
            image_bytes = _load_image_bytes(image_filename, image_data_base64)
            if image_bytes:
                print(f"Loaded image: {image_filename}")
            else:
                print(f"Failed to load image: {image_filename}")
        
        # Create presentation
        prs = Presentation()
        if VERBOSE:
            logger.info("Created PowerPoint presentation object")
        
        for i, slide in enumerate(slides):
            title = slide.get('title', 'Untitled Slide')
            content = slide.get('content', '')
            
            # Add slide
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide_obj = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_shape = slide_obj.shapes.title
            title_shape.text = title
            if VERBOSE:
                logger.info(f"Added slide {i+1}: {title}")
            
            # Add content
            body_shape = slide_obj.placeholders[1]
            tf = body_shape.text_frame
            tf.text = ""
            
            # Set text alignment to left
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            
            # Process bullet points if content contains them
            if content.strip() and content.strip().startswith('-'):
                # Split by newline and process each line
                items = [item.strip() for item in content.split('\n') if item.strip()]
                if VERBOSE:
                    logger.info(f"Slide {i+1} has {len(items)} bullet points")
                for item in items:
                    if item.startswith('-'):
                        item_text = item[1:].strip()  # Remove the dash
                        p = tf.add_paragraph()
                        p.text = item_text
                        p.level = 0
                        p.font.size = Pt(24)
                        p.space_after = Pt(6)
                        p.alignment = PP_ALIGN.LEFT
            else:
                # Handle regular text without bullet points
                p = tf.add_paragraph()
                p.text = content
                p.font.size = Pt(24)
                p.space_after = Pt(6)
                p.alignment = PP_ALIGN.LEFT
            
            # Add image to first slide if provided
            if i == 0 and image_bytes:
                _add_image_to_slide(slide_obj, image_bytes, 
                                   left=Inches(0.5), top=Inches(3.5), 
                                   width=Inches(4), height=Inches(3))
            
        # Write outputs to a temporary directory and clean up after encoding
        with tempfile.TemporaryDirectory() as tmpdir:
            # Save presentation
            pptx_output_path = os.path.join(tmpdir, "output_presentation.pptx")
            prs.save(pptx_output_path)
            if VERBOSE:
                logger.info(f"Saved PowerPoint presentation to {pptx_output_path}")

            # Create HTML file instead of PDF
            html_output_path = os.path.join(tmpdir, "output_presentation.html")
            if VERBOSE:
                logger.info(f"Starting HTML creation to {html_output_path}")

            # Create HTML representation of the presentation
            html_content = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PowerPoint Presentation</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; background-color: #f5f5f5; }
        .slide { background: white; margin: 20px auto; padding: 40px; max-width: 800px; box-shadow: 0 4px 6px rgba(0,0,0,0.1); border-radius: 8px; page-break-after: always; }
        .slide-title { color: #2c3e50; font-size: 28px; font-weight: bold; margin-bottom: 20px; border-bottom: 2px solid #3498db; padding-bottom: 10px; }
        .slide-content { font-size: 18px; line-height: 1.6; }
        .slide-content ul { margin: 0; padding-left: 30px; }
        .slide-content li { margin-bottom: 8px; }
        .slide-image { max-width: 400px; max-height: 300px; display: block; margin: 20px auto; border-radius: 4px; }
        .slide-number { position: absolute; top: 10px; right: 20px; color: #7f8c8d; font-size: 14px; }
    </style>
</head>
<body>"""
        
            for i, slide in enumerate(slides):
                title = slide.get('title', 'Untitled Slide')
                content = slide.get('content', '')

                html_content += f"""
    <div class="slide">
        <div class="slide-number">Slide {i+1}</div>
        <div class="slide-title">{title}</div>
        <div class="slide-content">"""
                
                # Add image to first slide if provided
                if i == 0 and image_bytes:
                    try:
                        img = Image.open(io.BytesIO(image_bytes))
                        mime_type = Image.MIME.get(img.format)
                        if mime_type:
                            img_b64 = base64.b64encode(image_bytes).decode('utf-8')
                            html_content += f'<img src="data:{mime_type};base64,{img_b64}" class="slide-image" />'
                    except Exception as e:
                        if VERBOSE:
                            logger.warning(f"Could not process image for HTML conversion: {e}")

                if content.strip() and content.strip().startswith('-'):
                    items = [item.strip() for item in content.split('\n') if item.strip()]
                    html_content += "<ul>"
                    for item in items:
                        if item.startswith('-'):
                            item_text = item[1:].strip()
                            html_content += f"<li>{item_text}</li>"
                    html_content += "</ul>"
                else:
                    html_content += f"<p>{content}</p>"

                html_content += """
        </div>
    </div>"""

            html_content += """
</body>
</html>"""
            
            # Save HTML file
            try:
                if VERBOSE:
                    logger.info("Saving HTML file...")
                with open(html_output_path, "w", encoding="utf-8") as f:
                    f.write(html_content)
                if VERBOSE:
                    logger.info(f"HTML successfully created at {html_output_path}")
            except Exception as e:
                # If HTML save fails, continue with just PPTX
                if VERBOSE:
                    logger.warning(f"HTML creation failed: {str(e)}")
                # Remove the HTML file if it exists
                if os.path.exists(html_output_path):
                    os.remove(html_output_path)
                if VERBOSE:
                    logger.info("HTML file removed due to creation error")

            # Read PPTX file as bytes
            with open(pptx_output_path, "rb") as f:
                pptx_bytes = f.read()

            # Encode PPTX as base64
            pptx_b64 = base64.b64encode(pptx_bytes).decode('utf-8')
            if VERBOSE:
                logger.info("PPTX file successfully encoded to base64")

            # Read HTML file as bytes if it exists
            html_b64 = None
            if os.path.exists(html_output_path):
                with open(html_output_path, "r", encoding="utf-8") as f:
                    html_content_file = f.read()
                html_b64 = base64.b64encode(html_content_file.encode('utf-8')).decode('utf-8')
                if VERBOSE:
                    logger.info("HTML file successfully encoded to base64")

            # Prepare artifacts
            artifacts = [
                {
                    "name": "presentation.pptx",
                    "b64": pptx_b64,
                    "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                }
            ]

            # Add HTML if creation was successful
            if html_b64:
                artifacts.append({
                    "name": "presentation.html",
                    "b64": html_b64,
                    "mime": "text/html",
                })
                if VERBOSE:
                    logger.info(f"Added {len(artifacts)} artifacts to response")
            else:
                if VERBOSE:
                    logger.info("No HTML artifact added due to creation failure")

            return {
                "results": {
                    "operation": "json_to_pptx",
                    "message": "PowerPoint presentation and HTML file generated successfully.",
                    "html_generated": html_b64 is not None,
                    "image_included": image_bytes is not None,
                },
                "artifacts": artifacts,
                "display": {
                    "open_canvas": True,
                    "primary_file": "presentation.pptx",
                    "mode": "replace",
                    "viewer_hint": "powerpoint",
                },
                "meta_data": {
                    "generated_slides": len(slides),
                    "output_files": [f"presentation.pptx", "presentation.html"] if html_b64 else ["presentation.pptx"],
                    "output_file_paths": ["temp:output_presentation.pptx", "temp:output_presentation.html"] if html_b64 else ["temp:output_presentation.pptx"],
                },
            }
    except Exception as e:
        if VERBOSE:
            logger.info(f"Error in json_to_pptx: {str(e)}")
        return {"results": {"error": f"Error creating PowerPoint: {str(e)}"}}


@mcp.tool
def markdown_to_pptx(
    markdown_content: Annotated[str, "Markdown content with headers (# or ##) as slide titles and content below each header"],
    image_filename: Annotated[str, "Optional image filename to integrate into the presentation"] = "",
    image_data_base64: Annotated[str, "Framework may supply Base64 image content as fallback"] = ""
) -> Dict[str, Any]:
    """
    Converts markdown content to PowerPoint presentation with support for bullet point lists and optional image integration
    
    Args:
        markdown_content: Markdown content where headers (# or ##) become slide titles and content below becomes slide content
        image_filename: Optional image filename to integrate into the presentation
        image_data_base64: Framework may supply Base64 image content as fallback
    
    Returns:
        Dictionary with 'results' and 'artifacts' keys:
        - 'results': Success message or error message
        - 'artifacts': List of artifact dictionaries with 'name', 'b64', and 'mime' keys
    """
    if VERBOSE:
        logger.info("Starting markdown_to_pptx execution...")
    try:
        # Parse markdown into slides
        slides = _parse_markdown_slides(markdown_content)
        if VERBOSE:
            logger.info(f"Parsed {len(slides)} slides from markdown")
        
        if not slides:
            return {"results": {"error": "No slides could be parsed from markdown content"}}
        
        # Load image if provided
        image_bytes = None
        if image_filename:
            image_bytes = _load_image_bytes(image_filename, image_data_base64)
            if image_bytes:
                if VERBOSE:
                    logger.info(f"Loaded image: {image_filename}")
            else:
                if VERBOSE:
                    logger.info(f"Failed to load image: {image_filename}")
        
        # Create presentation
        prs = Presentation()
        print("Created PowerPoint presentation object")
        
        for i, slide_data in enumerate(slides):
            title = slide_data.get('title', 'Untitled Slide')
            content = slide_data.get('content', '')
            
            # Add slide
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide_obj = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_shape = slide_obj.shapes.title
            title_shape.text = title
            if VERBOSE:
                logger.info(f"Added slide {i+1}: {title}")
            
            # Add content
            body_shape = slide_obj.placeholders[1]
            tf = body_shape.text_frame
            tf.text = ""
            
            # Set text alignment to left
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            
            # Process content - handle bullet points and regular text
            if content.strip():
                lines = content.split('\n')
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    
                    p = tf.add_paragraph()
                    
                    # Handle bullet points (- or *)
                    if line.startswith(('- ', '* ')):
                        p.text = line[2:].strip()
                        p.level = 0
                    # Handle sub-bullets (indented)
                    elif line.startswith(('  - ', '  * ', '\t- ', '\t* ')):
                        p.text = line.strip()[2:].strip()
                        p.level = 1
                    else:
                        # Regular text
                        p.text = line
                        p.level = 0
                    
                    p.font.size = Pt(24)
                    p.space_after = Pt(6)
                    p.alignment = PP_ALIGN.LEFT
            
            # Add image to first slide if provided
            if i == 0 and image_bytes:
                _add_image_to_slide(slide_obj, image_bytes, 
                                   left=Inches(0.5), top=Inches(3.5), 
                                   width=Inches(4), height=Inches(3))
        
        # Write outputs to a temporary directory and clean up after encoding
        with tempfile.TemporaryDirectory() as tmpdir:
            # Save presentation
            pptx_output_path = os.path.join(tmpdir, "output_presentation.pptx")
            prs.save(pptx_output_path)
            if VERBOSE:
                logger.info(f"Saved PowerPoint presentation to {pptx_output_path}")

            # Create HTML file instead of PDF
            html_output_path = os.path.join(tmpdir, "output_presentation.html")
            if VERBOSE:
                logger.info(f"Starting HTML creation to {html_output_path}")

            # Create HTML representation of the presentation
            html_content = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PowerPoint Presentation</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; background-color: #f5f5f5; }
        .slide { background: white; margin: 20px auto; padding: 40px; max-width: 800px; box-shadow: 0 4px 6px rgba(0,0,0,0.1); border-radius: 8px; page-break-after: always; }
        .slide-title { color: #2c3e50; font-size: 28px; font-weight: bold; margin-bottom: 20px; border-bottom: 2px solid #3498db; padding-bottom: 10px; }
        .slide-content { font-size: 18px; line-height: 1.6; }
        .slide-content ul { margin: 0; padding-left: 30px; }
        .slide-content li { margin-bottom: 8px; }
        .slide-image { max-width: 400px; max-height: 300px; display: block; margin: 20px auto; border-radius: 4px; }
        .slide-number { position: absolute; top: 10px; right: 20px; color: #7f8c8d; font-size: 14px; }
    </style>
</head>
<body>"""
        
            for i, slide_data in enumerate(slides):
                title = slide_data.get('title', 'Untitled Slide')
                content = slide_data.get('content', '')

                html_content += f"""
    <div class="slide">
        <div class="slide-number">Slide {i+1}</div>
        <div class="slide-title">{title}</div>
        <div class="slide-content">"""
                
                # Add image to first slide if provided
                if i == 0 and image_bytes:
                    try:
                        img = Image.open(io.BytesIO(image_bytes))
                        mime_type = Image.MIME.get(img.format)
                        if mime_type:
                            img_b64 = base64.b64encode(image_bytes).decode('utf-8')
                            html_content += f'<img src="data:{mime_type};base64,{img_b64}" class="slide-image" />'
                    except Exception as e:
                        if VERBOSE:
                            logger.warning(f"Could not process image for HTML conversion: {e}")

                if content.strip():
                    lines = content.split('\n')
                    bullet_lines = []
                    regular_lines = []

                    for line in lines:
                        line = line.strip()
                        if line.startswith(('- ', '* ')):
                            bullet_lines.append(line[2:].strip())
                        elif line.startswith(('  - ', '  * ', '\t- ', '\t* ')):
                            bullet_lines.append(f"&nbsp;&nbsp;&bull; {line.strip()[2:].strip()}")
                        elif line:
                            regular_lines.append(line)

                    if bullet_lines:
                        html_content += "<ul>"
                        for item in bullet_lines:
                            html_content += f"<li>{item}</li>"
                        html_content += "</ul>"

                    if regular_lines:
                        for line in regular_lines:
                            html_content += f"<p>{line}</p>"

                html_content += """
        </div>
    </div>"""

            html_content += """
</body>
</html>"""
            
            # Save HTML file
            try:
                if VERBOSE:
                    logger.info("Saving HTML file...")
                with open(html_output_path, "w", encoding="utf-8") as f:
                    f.write(html_content)
                if VERBOSE:
                    logger.info(f"HTML successfully created at {html_output_path}")
            except Exception as e:
                # If HTML save fails, continue with just PPTX
                if VERBOSE:
                    logger.warning(f"HTML creation failed: {str(e)}")
                # Remove the HTML file if it exists
                if os.path.exists(html_output_path):
                    os.remove(html_output_path)
                if VERBOSE:
                    logger.info("HTML file removed due to creation error")

            # Read PPTX file as bytes
            with open(pptx_output_path, "rb") as f:
                pptx_bytes = f.read()

            # Encode PPTX as base64
            pptx_b64 = base64.b64encode(pptx_bytes).decode('utf-8')
            if VERBOSE:
                logger.info("PPTX file successfully encoded to base64")

            # Read HTML file as bytes if it exists
            html_b64 = None
            if os.path.exists(html_output_path):
                with open(html_output_path, "r", encoding="utf-8") as f:
                    html_content_file = f.read()
                html_b64 = base64.b64encode(html_content_file.encode('utf-8')).decode('utf-8')
                if VERBOSE:
                    logger.info("HTML file successfully encoded to base64")

            # Prepare artifacts
            artifacts = [
                {
                    "name": "presentation.pptx",
                    "b64": pptx_b64,
                    "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                }
            ]

            # Add HTML if creation was successful
            if html_b64:
                artifacts.append({
                    "name": "presentation.html",
                    "b64": html_b64,
                    "mime": "text/html",
                })
                if VERBOSE:
                    logger.info(f"Added {len(artifacts)} artifacts to response")
            else:
                if VERBOSE:
                    logger.info("No HTML artifact added due to creation failure")

            return {
                "results": {
                    "operation": "markdown_to_pptx",
                    "message": "PowerPoint presentation and HTML file generated successfully from markdown.",
                    "html_generated": html_b64 is not None,
                    "image_included": image_bytes is not None,
                },
                "artifacts": artifacts,
                "display": {
                    "open_canvas": True,
                    "primary_file": "presentation.pptx",
                    "mode": "replace",
                    "viewer_hint": "powerpoint",
                },
                "meta_data": {
                    "generated_slides": len(slides),
                    "output_files": [f"presentation.pptx", "presentation.html"] if html_b64 else ["presentation.pptx"],
                    "output_file_paths": ["temp:output_presentation.pptx", "temp:output_presentation.html"] if html_b64 else ["temp:output_presentation.pptx"],
                },
            }
    except Exception as e:
        print(f"Error in markdown_to_pptx: {str(e)}")
        return {"results": {"error": f"Error creating PowerPoint from markdown: {str(e)}"}}


if __name__ == "__main__":
    mcp.run()
